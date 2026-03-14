"""Extract text content from Word, PowerPoint, and Excel files."""

import io
import logging
from pathlib import Path

from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".xlsx"}


def parse_document(filename: str, content: bytes) -> str:
    """Parse a document and return extracted text.

    Raises ValueError for unsupported file types.
    """
    ext = Path(filename).suffix.lower()

    if ext == ".docx":
        return _parse_docx(content)
    elif ext == ".pptx":
        return _parse_pptx(content)
    elif ext == ".xlsx":
        return _parse_xlsx(content)
    else:
        raise ValueError(
            f"Unsupported file type: {ext}. "
            f"Supported types: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )


def _parse_docx(content: bytes) -> str:
    doc = DocxDocument(io.BytesIO(content))
    parts: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # Preserve heading structure
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            level = para.style.name.replace("Heading", "").strip()
            prefix = "#" * int(level) if level.isdigit() else "##"
            parts.append(f"{prefix} {text}")
        else:
            parts.append(text)

    # Also extract text from tables
    for table in doc.tables:
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            # First row as header
            parts.append(rows[0])
            parts.append(" | ".join(["---"] * len(table.rows[0].cells)))
            parts.extend(rows[1:])

    return "\n\n".join(parts)


def _parse_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_parts.append(" | ".join(cells))

        if slide_parts:
            parts.append(f"## Slide {i}\n" + "\n".join(slide_parts))

    return "\n\n".join(parts)


def _parse_xlsx(content: bytes) -> str:
    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):  # skip fully empty rows
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"## Sheet: {sheet_name}\n" + "\n".join(rows))

    wb.close()
    return "\n\n".join(parts)
