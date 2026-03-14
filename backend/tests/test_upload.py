import io

import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches

from app.services.sessions import session_store


def _make_docx(text: str = "Hello world") -> bytes:
    doc = DocxDocument()
    doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx(text: str = "Slide content") -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.mark.anyio
async def test_upload_docx(client):
    session_id = (await client.post("/api/chat/sessions")).json()["session_id"]

    content = _make_docx("Strategic initiative overview")
    resp = await client.post(
        f"/api/chat/upload?session_id={session_id}",
        files=[("files", ("report.docx", content, "application/vnd.openxmlformats-officedocument.wordprocessingml.document"))],
    )
    assert resp.status_code == 200
    results = resp.json()
    assert len(results) == 1
    assert results[0]["filename"] == "report.docx"
    assert results[0]["ok"] is True

    # Verify document is stored in session
    docs = session_store.get_documents(session_id)
    assert len(docs) == 1
    assert "Strategic initiative overview" in docs[0].text


@pytest.mark.anyio
async def test_upload_pptx(client):
    session_id = (await client.post("/api/chat/sessions")).json()["session_id"]

    content = _make_pptx("Key findings")
    resp = await client.post(
        f"/api/chat/upload?session_id={session_id}",
        files=[("files", ("deck.pptx", content, "application/vnd.openxmlformats-officedocument.presentationml.presentation"))],
    )
    assert resp.status_code == 200
    results = resp.json()
    assert results[0]["ok"] is True

    docs = session_store.get_documents(session_id)
    assert "Key findings" in docs[0].text


@pytest.mark.anyio
async def test_upload_multiple_files(client):
    session_id = (await client.post("/api/chat/sessions")).json()["session_id"]

    docx = _make_docx("Document one")
    pptx = _make_pptx("Document two")
    resp = await client.post(
        f"/api/chat/upload?session_id={session_id}",
        files=[
            ("files", ("a.docx", docx, "application/octet-stream")),
            ("files", ("b.pptx", pptx, "application/octet-stream")),
        ],
    )
    assert resp.status_code == 200
    results = resp.json()
    assert len(results) == 2
    assert all(r["ok"] for r in results)

    docs = session_store.get_documents(session_id)
    assert len(docs) == 2


@pytest.mark.anyio
async def test_upload_unsupported_file(client):
    session_id = (await client.post("/api/chat/sessions")).json()["session_id"]

    resp = await client.post(
        f"/api/chat/upload?session_id={session_id}",
        files=[("files", ("notes.txt", b"hello", "text/plain"))],
    )
    assert resp.status_code == 200
    results = resp.json()
    assert len(results) == 1
    assert results[0]["ok"] is False
    assert "Unsupported" in results[0]["error"]


@pytest.mark.anyio
async def test_upload_invalid_session(client):
    resp = await client.post(
        "/api/chat/upload?session_id=nonexistent",
        files=[("files", ("a.docx", _make_docx(), "application/octet-stream"))],
    )
    assert resp.status_code == 404


@pytest.mark.anyio
async def test_upload_then_chat_includes_context(client):
    """Documents uploaded to a session should be available as context in chat."""
    session_id = (await client.post("/api/chat/sessions")).json()["session_id"]

    # Upload a document
    content = _make_docx("Revenue grew 25% year over year")
    await client.post(
        f"/api/chat/upload?session_id={session_id}",
        files=[("files", ("report.docx", content, "application/octet-stream"))],
    )

    # Chat in that session — the agent should have document context
    resp = await client.post(
        "/api/chat",
        json={"message": "What does the report say?", "session_id": session_id},
    )
    assert resp.status_code == 200
    # We can't check the agent actually used the context (TestModel),
    # but we verify the flow doesn't error
    assert resp.json()["message"] == "Test response"
