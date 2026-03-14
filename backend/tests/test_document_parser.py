import io

import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from app.services.document_parser import parse_document


def _make_docx(paragraphs: list[tuple[str, str | None]]) -> bytes:
    """Create a .docx in memory. Each item is (text, style_name_or_None)."""
    doc = DocxDocument()
    for text, style in paragraphs:
        doc.add_paragraph(text, style=style)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_docx_with_table(rows: list[list[str]]) -> bytes:
    doc = DocxDocument()
    doc.add_paragraph("Before table")
    table = doc.add_table(rows=len(rows), cols=len(rows[0]))
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            table.rows[i].cells[j].text = cell_text
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx(slides: list[list[str]]) -> bytes:
    """Create a .pptx with one text box per string per slide."""
    prs = Presentation()
    for texts in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        for i, text in enumerate(texts):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(5), Inches(1))
            txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx(sheets: dict[str, list[list[str]]]) -> bytes:
    wb = Workbook()
    first = True
    for name, rows in sheets.items():
        ws = wb.active if first else wb.create_sheet(name)
        if first:
            ws.title = name
            first = False
        for row in rows:
            ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# --- docx tests ---

def test_parse_docx_paragraphs():
    content = _make_docx([
        ("Introduction", "Heading 1"),
        ("This is body text.", None),
        ("Details", "Heading 2"),
        ("More body text.", None),
    ])
    result = parse_document("test.docx", content)
    assert "# Introduction" in result
    assert "## Details" in result
    assert "This is body text." in result
    assert "More body text." in result


def test_parse_docx_table():
    content = _make_docx_with_table([
        ["Name", "Score"],
        ["Alice", "95"],
        ["Bob", "87"],
    ])
    result = parse_document("test.docx", content)
    assert "Before table" in result
    assert "Name | Score" in result
    assert "Alice | 95" in result


def test_parse_docx_empty():
    content = _make_docx([])
    result = parse_document("empty.docx", content)
    assert result == ""


# --- pptx tests ---

def test_parse_pptx_basic():
    content = _make_pptx([
        ["Title Slide", "Subtitle"],
        ["Second slide content"],
    ])
    result = parse_document("deck.pptx", content)
    assert "## Slide 1" in result
    assert "Title Slide" in result
    assert "Subtitle" in result
    assert "## Slide 2" in result
    assert "Second slide content" in result


def test_parse_pptx_empty():
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank slide, no text
    buf = io.BytesIO()
    prs.save(buf)
    result = parse_document("empty.pptx", buf.getvalue())
    assert result == ""


# --- xlsx tests ---

def test_parse_xlsx_basic():
    content = _make_xlsx({
        "Revenue": [["Q1", "Q2", "Q3"], ["100", "200", "300"]],
    })
    result = parse_document("data.xlsx", content)
    assert "## Sheet: Revenue" in result
    assert "Q1 | Q2 | Q3" in result
    assert "100 | 200 | 300" in result


def test_parse_xlsx_multiple_sheets():
    content = _make_xlsx({
        "Sheet1": [["A", "B"]],
        "Sheet2": [["C", "D"]],
    })
    result = parse_document("multi.xlsx", content)
    assert "## Sheet: Sheet1" in result
    assert "## Sheet: Sheet2" in result


# --- error handling ---

def test_unsupported_extension():
    with pytest.raises(ValueError, match="Unsupported file type"):
        parse_document("file.pdf", b"fake content")


def test_unsupported_extension_txt():
    with pytest.raises(ValueError, match="Unsupported file type"):
        parse_document("notes.txt", b"hello")
